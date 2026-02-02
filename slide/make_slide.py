import os
import re
import json
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm  # プログレスバー

# 🔹 設定ファイル
PDF_FILE = "slide.pdf"       # 入力PDF
NOTES_FILE = "notes.json"     # スライドノート用JSON
IMAGE_DIR = "slide"           # 画像保存フォルダ
OUTPUT_FILE = "slide.pptx"   # 出力PowerPointファイル

# PowerPointスライドサイズ設定（16:9）
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# 🔹 1️⃣ PDF → PNG変換（プログレスバー付き）
if not os.path.exists(IMAGE_DIR):
    os.makedirs(IMAGE_DIR)

print("📄 PDFをPNGに変換中...")
images = convert_from_path(PDF_FILE, dpi=300)  # 300dpiで高画質変換
for i, img in tqdm(enumerate(images), total=len(images), desc="Converting PDF"):
    img_path = os.path.join(IMAGE_DIR, f"slide-{i+1:02d}.png")
    img.save(img_path, "PNG")

# 🔹 2️⃣ ノートJSONを読み込み
notes_data = {}
if os.path.exists(NOTES_FILE):
    with open(NOTES_FILE, "r", encoding="utf-8") as f:
        notes_data = json.load(f)

# 🔹 3️⃣ PowerPointに画像とノートを追加（プログレスバー付き）
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

slide_pattern = re.compile(r"slide-(\d{2})\.png")
image_files = [f for f in os.listdir(IMAGE_DIR) if slide_pattern.match(f)]
image_files.sort()

print("📊 PowerPointスライドを作成中...")
for filename in tqdm(image_files, total=len(image_files), desc="Creating PPTX"):
    slide_num = int(filename.split("-")[1].split(".")[0])  # スライド番号を取得
    img_path = os.path.join(IMAGE_DIR, filename)

    # スライド追加
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 画像を左上に配置し、スライドサイズいっぱいに拡大
    picture = slide.shapes.add_picture(img_path, 0, 0)
    picture.width = SLIDE_WIDTH
    picture.height = SLIDE_HEIGHT

    # 🔹 スライドノートを追加
    if str(slide_num) in notes_data:
        slide.notes_slide.notes_text_frame.text = notes_data[str(slide_num)]

# 🔹 4️⃣ PowerPointファイルを保存
prs.save(OUTPUT_FILE)
print(f"✅ PowerPointファイル '{OUTPUT_FILE}' が作成されました！🎉")
